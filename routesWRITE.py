import sys, boto3, random, base64, os, secrets, httplib2, json, ast
from sqlalchemy import asc, desc
from datetime import datetime, timedelta, date
from flask import render_template, url_for, flash, redirect, request, abort, jsonify
from app import app, db, bcrypt, mail
from flask_login import login_user, current_user, logout_user, login_required
from forms import *
from models import *
from routesGet import getUsers
from pprint import pprint
from flask_mail import Message

from meta import BaseConfig
s3_resource = BaseConfig.s3_resource

""" improvements """

# limit word count for plans
## seperate plan into idea + 2 details  with word limit 6
# minimum sentences for writing
# maximum word count for parts - stop monster eassy
# record if pasting has occurred https://stackoverflow.com/questions/3211505/detect-pasted-text-with-ctrlv-or-right-click-paste
# monitor titles (all caps, no caps, periods)
# figure out why when students fix title it doesn't update
# Why is draft sometimes saving as stage 1?
## Do not allow work to be saved if 'i' or no space between . ,
### Do not allow work to be save if first letter is not capital

## allow instructor to edit plan from dash

## purple revise color should include capital letters too
# find a way to open units for students more easily
# figure out why the stage is recorded in the josn and the grade column


"""S3 CONNECTIONS """
def loadAWS(file, unit):
    SCHEMA = getSchema()
    S3_LOCATION = schemaList[SCHEMA]['S3_LOCATION']
    S3_BUCKET_NAME = schemaList[SCHEMA]['S3_BUCKET_NAME']

    if unit == 0 :
        key = file
        content_object = s3_resource.Object(  S3_BUCKET_NAME, key  )
        print('LOAD_AWS_SEARCH', content_object)
    else:
        key = str(unit) + '/' + file
        content_object = s3_resource.Object(  S3_BUCKET_NAME, key   )
        print('LOAD_AWS_SEARCH', content_object)
    try:
        file_content = content_object.get()['Body'].read().decode('utf-8')
        jload = json.loads(file_content)
        print(type(jload))
    except:
        print('AWS return NONE')
        jload = None

    return jload

# def loadAWS(file, unit):
#     if unit == 0 :
#         jList = [S3_BUCKET_NAME, file]
#         content_object = s3_resource.Object(  jList[0], jList[1]   )
#     else:
#         jList = [S3_BUCKET_NAME, file]
#         content_object = s3_resource.Object(  jList[0], jList[1]   )

#     file_content = content_object.get()['Body'].read().decode('utf-8')
#     jload = json.loads(file_content)
#     print(type(jload))

#     return jload


def getWriteUsers():
    ### sperate for vietnam class

    vs = []
    es = []
    userList = []

    ## CODE FOR SEPARTING TO CLASSES

    # for u in getUsers(10):
    #     # print(u.username, u.vtm)
    #     vs.append(u.username)

    # if current_user.username not in vs:

    #     for u in getUsers(8):
    #         if u.username not in vs:
    #             es.append(u.username)

    # if current_user.username not in vs:
    #     userList = es
    # else:
    #     userList = vs

    for u in getUsers(10):
        es.append(u.username)


    userList = es

    return userList

def get_all_values(nested_dictionary):
    detected = 0
    for key, value in nested_dictionary.items():
        if type(value) is dict:
            print ('DICT FOUND', value)
            if get_all_values(value) != 0:
                detected += get_all_values(value)
        else:
            if value == None or value == "":
                print('CHECK', key, value)
                detected += 1

    return detected

@app.route('/updateWritingPresentation', methods=['POST'])
def updateWritingPresentation():

    projectData = {
        1 : U011U_WRITE,
        2 : U012U_WRITE
    }

    setup = int(request.form ['proj'])
    ansOBJ = request.form ['ansOBJ']

    ansDict = json.loads(ansOBJ)
    print(setup, ansDict)
    try:
        ## instructor update
        grade = request.form ['grade']
        name = request.form ['name']
    except:
        grade = int(request.form ['grade'])
        name = current_user.username

        if grade == 5:
            grade = 5
        elif get_all_values(ansDict) == 0:
            grade = 1
        else:
            print ('GET_ALL', get_all_values(ansDict))
            grade = 0

    print('GRADE', grade)

    project_answers = projectData[setup].query.filter_by(username=name).first()
    project_answers.Ans01 = json.dumps(ansDict)

    project_answers.Grade = grade
    db.session.commit()

    return jsonify({'grade' : grade})


@app.route ("/write_projects", methods=['GET','POST'])
@login_required
def write_projects():
    SCHEMA = getSchema()

    users = getUsers(SCHEMA)

    mtDict = {}
    for user in users:
        mtDict[user.username] = {
            'ID' : user.studentID,
            'Grade' : 0,
            'Data' : {}
        }

    project_answers = U011U_WRITE.query.all()

    for answer in project_answers:
        proj = None
        projDict = json.loads(answer.Ans01)

        # print(get_all_values(cvDict))

        if answer.Grade == 0:
            tGrade = 1
        else:
            tGrade = answer.Grade

        mtDict[answer.username]['Proj'] = proj
        mtDict[answer.username]['Grade'] = tGrade
        mtDict[answer.username]['Data'] = projDict



    pprint(mtDict)
    source = 'https://docs.google.com/presentation/d/e/2PACX-1vSX1RPIa4phQJsdYa-siNiC2W_by9LL_xKBMGsyLkFAFnN2eaCd7w3k9Svd9rv64J2c6u9NSSZMP6PS/embed'




    return render_template('work/write_projects.html', legend='Presentation Project',
    source1=source, source2='', ansString=json.dumps(mtDict)  )


@app.route('/createPPTwrite', methods=['POST'])
def createPPTwrite():
    SCHEMA = getSchema()
    S3_LOCATION = schemaList[SCHEMA]['S3_LOCATION']
    S3_BUCKET_NAME = schemaList[SCHEMA]['S3_BUCKET_NAME']

    proj = request.form ['proj']
    ansOBJ = request.form ['ansOBJ']

    ansDict = json.loads(ansOBJ)
    print(ansDict)

    head = ansDict['Title'].title()

    if get_all_values(ansDict) != 0:
        print('ERROR')
        return jsonify({'error' : 100})

    from pptx import Presentation

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Workplace Presentation"
    subtitle.text = head + "\nPresentation by "


    intro = {
        'Situation' : ansDict['situation_kw'],
        'Problem' : ansDict['problem_kw'],
        'Solution' : ansDict['solution_kw']
    }

    for i in intro:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = i
        tf = body_shape.text_frame
        kw = intro[i].split(',')
        print(kw)
        for ckw in kw:
            p = tf.add_paragraph()
            p.text = ckw.title().strip()
            p.level = 1

    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = 'Reasons'
    for r in ansDict['Reasons']:
        p = tf.add_paragraph()
        p.text = ansDict['Reasons'][r]
        p.level = 1

    count = 1
    for part in ansDict['Parts']:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title

        title_shape.text = str(count) + ': ' + ansDict['Reasons'][part]

        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        p_1 = tf.add_paragraph()
        p_1.text = '1'
        p_1.level = 1
        idea_kw = ansDict['Parts'][part]['kw']['1'].split(',')
        for i in idea_kw:
            p = tf.add_paragraph()
            p.text = i.title().strip()
            p.level = 2

        p_2 = tf.add_paragraph()
        p_2.text = '2'
        p_2.level = 1

        idea_kw_2 = ansDict['Parts'][part]['kw']['2'].split(',')
        for i in idea_kw_2:
            p = tf.add_paragraph()
            p.text = i.title().strip()
            p.level = 2

        p_3 = tf.add_paragraph()
        p_3.text = '3'
        p_3.level = 1
        idea_kw_3 = ansDict['Parts'][part]['kw']['3'].split(',')
        for i in idea_kw_3:
            p = tf.add_paragraph()
            p.text = i.title().strip()
            p.level = 2

        count +=1

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Conclusion'
    tf = body_shape.text_frame
    closing_kw = ansDict['closing_kw'].split(',')
    print(closing_kw)
    for ckw in closing_kw:
        p = tf.add_paragraph()
        p.text = ckw.strip().title()
        p.level = 1

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    script_1 = ['Welcome to our presentation, today we would like to discuss the topic of ____________',
              "Let's introduce the situation",
              ansDict['Situation'],
              ansDict['Problem'],
              ansDict['Solution']
            ]
    for part in script_1:
        p = tf.add_paragraph()
        p.text = part
        p.level = 1

    for part in ansDict['Parts']:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title

        title_shape.text = str(count) + ': ' + ansDict['Reasons'][part]

        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        p = tf.add_paragraph()
        idea_dt = ansDict['Parts'][part]['dt']
        for i in idea_dt:
            p = tf.add_paragraph()
            p.text = idea_dt[i]
            p.level = 2


    print('PROCESSING PPT')
    filename = current_user.username + proj + '.pptx'
    try:
        os.remove(filename)
    except:
        print('No OS File')
    prs.save(filename)
    data = open(filename, 'rb')
    aws_filename = 'projects/' + filename
    pptLink = S3_LOCATION + aws_filename

    s3_resource.Bucket(S3_BUCKET_NAME).put_object(Key=aws_filename, Body=data)
    print(filename)

    return jsonify({'pptLink' : pptLink})


@app.route ("/work_presentation_list", methods=['GET','POST'])
@login_required
def work_presentation_list():


    source = 'https://docs.google.com/presentation/d/e/2PACX-1vQmzC3mSJaWwaSZxN2O0iOvelYQCHmaoeH2cCTflVU4MJoDwAH6u6bshtzkIJpsXkEarJEKGkRVZf64/embed'
    setup = 1

    projectData = {
        1 : U011U_WRITE,
        2 : U012U_WRITE
    }

    startDict = {
            'Title' : None,
            'Situation' : None,
            'situation_kw' : None,
            'Problem' : None,
            'problem_kw' : None,
            'Solution' : None,
            'solution_kw' : None,
            'Closing' : None,
            'closing_kw' : None,
            'Reasons' : {
                1 : None,
                2 : None,
                3 : None
            },
            'Parts' : {
                1 : {
                    'kw' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                    'dt' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                },
                2 : {
                    'kw' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                    'dt' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                },
                3 : {
                    'kw' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                    'dt' : {
                        1 : None,
                        2 : None,
                        3 : None,
                    },
                },
            }
        }


    if projectData[setup].query.filter_by(username=current_user.username).first():
        pass
    else:
        start = projectData[setup](username=current_user.username, Ans01=json.dumps(startDict), Grade=0, Comment='0')
        db.session.add(start)
        db.session.commit()


    project = projectData[setup].query.filter_by(username=current_user.username).first()
    ansDict = project.Ans01
    grade = project.Grade
    stage = project.Comment
    print(source)

    return render_template('work/work_list.html', legend='Presentation Projects', source=source, stage=stage, grade=grade, setup=setup)


@app.route ("/work_presentation_1", methods=['GET','POST'])
@login_required
def work_presentation_1():


    source = 'https://docs.google.com/presentation/d/e/2PACX-1vSX1RPIa4phQJsdYa-siNiC2W_by9LL_xKBMGsyLkFAFnN2eaCd7w3k9Svd9rv64J2c6u9NSSZMP6PS/embed'

    project = U011U_WRITE.query.filter_by(username=current_user.username).first()
    ansString = project.Ans01
    grade = project.Grade
    stage = project.Comment
    print(ansString)
    print(source)

    return render_template('work/work_presentation_1.html', legend='Presentation Projects', ansString=ansString, source=source, stage=stage, grade=grade)






@app.route("/tips", methods = ['GET', 'POST'])
@login_required
def tips():

    with open('static/topics.json', 'r') as f:
        srcJSON = json.load(f)

    tips = json.dumps(srcJSON['tips'])

    return render_template('work/tips.html', tips=tips)


@app.route('/upload/<string:assignment>', methods=['POST', 'GET'])
def upload(assignment):
    SCHEMA = getSchema()
    S3_LOCATION = schemaList[SCHEMA]['S3_LOCATION']
    S3_BUCKET_NAME = schemaList[SCHEMA]['S3_BUCKET_NAME']

    file = request.files['file']
    fn = file.filename
    file_name = current_user.username + '_' + fn + '_' + assignment + '.mp3'
    my_bucket = s3_resource.Bucket(S3_BUCKET_NAME)
    my_bucket.Object(file_name).put(Body=file)

    return redirect(request.referrer)


''' new style '''
@app.route('/storeData', methods=['POST'])
def storeData():
    unit = request.form ['unit']
    obj = request.form ['obj']
    stage = request.form ['stage']
    work = request.form ['work']

    if work == 'edit':
        student = request.form ['student']
    else:
        student = current_user.username

    print('STAGE WORK OBJ', stage, work, obj)
    classModel = getInfo()['aModsDict'][unit]
    entry = classModel.query.filter_by(username=student).first()
    info = json.loads(entry.info)

    if work == 'plan':
        if int(stage) == 0:
            info[work + "_date_start"] = str(date.today())
        if int(stage) == 1:
            info[work + "_date_finish"] = str(date.today())
            info['stage'] = 1
            entry.grade = stage
        entry.info = json.dumps(info)
        entry.plan = obj
        db.session.commit()

    if work == 'draft':
        if int(stage) == 1:
            info[work + "_date_start"] = str(date.today())
        if int(stage) == 2:
            info[work + "_date_finish"] = str(date.today())
            info['stage'] = 2
        entry.info = json.dumps(info)
        entry.draft = obj
        entry.grade = stage
        db.session.commit()

    if work == 'edit':
        dataDict = {
            'html' : request.form ['html'],
            'text' : request.form ['text'],
            'revised' : None,
        }
        info['stage'] = stage
        entry.info = json.dumps(info)
        entry.revise = json.dumps(dataDict)
        entry.grade = 3
        db.session.commit()
        studentEmail = User.query.filter_by(username=student).first().email
        msg = Message('Message from Writing LMS',
                sender='chrisflask0212@gmail.com',
                recipients=['cjx02121981@gmail.com', studentEmail ])
        msg.body = 'Dear ' + student + ',  Your writing draft for topic ' + unit + ' has been checked. Please move to the revise stage to see the correction and fix them before publishing. Thanks, Chris'
        mail.send(msg)

    if work == 'revise':
        print(student)
        info[work + "_date_finish"] = str(date.today())

        if int(info['stage']) == 3:
            info['stage'] = 4
            entry.grade = 4
        entry.info = json.dumps(info)
        print(entry.info)
        dataDict = json.loads(entry.revise)
        dataDict['revised'] = obj ## actually a string
        entry.revise = json.dumps(dataDict)
        print(entry.revise)
        db.session.commit()

    if work == 'publish':
        info[work + "_date_finish"] = str(date.today())
        info['stage'] = stage
        entry.info = json.dumps(info)
        dataDict = {
            'title' :  request.form ['title'],
            'imageLink' :  request.form ['imageLink'],
            'final' :  request.form ['final'],
        }
        entry.publish = json.dumps(dataDict)
        entry.grade = 5
        db.session.commit()


    name = current_user.username

    return jsonify({'name' : name, 'work' : work})
''' new style '''

@app.route('/updateInfo', methods=['POST'])
def updateInfo():

    unit = request.form ['unit']
    theme = request.form ['theme']
    name = request.form ['name']
    avatar = request.form ['avatar']
    partner = request.form ['partner']

    partner = partner.strip().title()

    print(avatar, partner, getWriteUsers())

    if partner == current_user.username:
        partner = 'None'
    elif partner and len(partner) > 0 and partner not in getWriteUsers():
        return jsonify({'partner' : partner})
    elif partner and partner == current_user.username:
        return jsonify({'partner' : partner})

    if len(name) > 0 and name in getWriteUsers():
        print('AVATAR NAME', len(name), name in getWriteUsers())
        return jsonify({'name' : name})

    print('UInfo', unit, theme, partner)
    classModel = getInfo()['aModsDict'][unit]
    entry = classModel.query.filter_by(username=current_user.username).first()

    info = json.loads(entry.info)

    info['theme'] = theme
    info['name'] = name
    info['avatar'] = avatar

    entry.info = json.dumps(info)
    entry.partner = partner

    db.session.commit()

    return jsonify({'success' : True})

@app.route('/switchPartner', methods=['POST'])
def switchPartner():

    unit = request.form ['unit']
    partner = request.form ['partner']

    partner = partner.strip().title()


    classModel = getInfo()['aModsDict'][unit]
    entry = classModel.query.filter_by(partner=partner).first()
    writer = entry.username

    print('U switch Info', unit, partner, writer)

    entry.partner = writer
    entry.username = partner
    db.session.commit()

    return jsonify({'success' : unit})


@app.route('/sendImage', methods=['POST'])
def sendImage():
    SCHEMA = getSchema()
    S3_LOCATION = schemaList[SCHEMA]['S3_LOCATION']
    S3_BUCKET_NAME = schemaList[SCHEMA]['S3_BUCKET_NAME']

    b64String = request.form ['b64String']
    print(b64String)
    print('SENDIMAGE ACTIVE')
    unit = request.form ['unit']
    fileType = request.form ['fileType']
    print (b64String, fileType)

    image = base64.b64decode(b64String)
    imageLink = S3_LOCATION + 'images/' + unit + '/' + current_user.username + '.' + fileType
    filename = 'images/' + unit + '/' + current_user.username + '.' + fileType
    s3_resource.Bucket(S3_BUCKET_NAME).put_object(Key=filename, Body=image)

    return jsonify({'name' : current_user.username, 'imageLink' : imageLink})

""" ### TOPICS ### """

## to do list
## controls image
## limited access to units

@app.route("/topic_list", methods = ['GET', 'POST'])
@login_required
def topic_list():
    topDict = {}
    infoDict = {}

    with open('static/topics.json', 'r') as f:
        srcJSON = json.load(f)

    sources = srcJSON['sources']

    # pprint(sources)

    userList = getWriteUsers()


    for unit in sources:
        src = sources[unit]
        if src['Set'] >= 1:
            ## add to topics list
            topDict[unit] = {
                'Title' : src['Title'],
                'Deadline' : src['Deadline'],
                'Start' : src['Date'],
                'Set' : src['Set']
            }
            model = getInfo()['aModsDict'][unit]

            user = None

            userW = model.query.filter_by(username=current_user.username).first()
            userP = model.query.filter_by(partner=current_user.username).first()

            print('source set', userW, userP)

            if userP and userW:
                planW = json.loads(userW.plan)
                planP = json.loads(userP.plan)
                print('double set', planW, planP)
                if planW == {}:
                    db.session.delete(userW)
                    db.session.commit()
                    user = userP
            elif userP:
                user = userP
            elif userW:
                user = userW



            ## add info details if available
            if user:
                infoDict = json.loads(user.info)
                topDict[unit]['Theme'] = infoDict['theme']
                topDict[unit]['Stage'] = int(infoDict['stage'])
                topDict[unit]['Avatar'] = infoDict['avatar']
                topDict[unit]['Partner'] = user.partner
                topDict[unit]['Writer'] = user.username
            else:
                topDict[unit]['Theme'] = 'white'
                topDict[unit]['Stage'] = 0
                topDict[unit]['Avatar'] = 'none'
                topDict[unit]['Writer'] = None

    print('DICT', topDict)
    topJS = json.dumps(topDict)
    topInfo = json.dumps(infoDict)

    return render_template('work/topic_list.html', topJS=topJS, topInfo=topInfo)


## AJAX
@app.route('/topicCheck/<string:unit>', methods=['POST'])
def topicCheck(unit):
    print('TOPIC CHECK ACTIVE')

    stage = 0
    dataList =  []

    userList = getWriteUsers()

    model = getInfo()['aModsDict'][unit]
    entryUser = model.query.filter_by(username=current_user.username).first()
    if not entryUser:
        entryUser = model.query.filter_by(partner=current_user.username).first()

    entries = model.query.all()

    for entry in entries:
        info = json.loads(entry.info)
        if entry.username == current_user.username or entry.partner == current_user.username:
            stage = info['stage']
            print('CURRENT USER FOUND', current_user.username, stage)
        elif entry.username in userList:
            plan = json.loads(entry.plan)
            draft = json.loads(entry.draft)
            publish = json.loads(entry.publish)

            entryDict = {
                'info' : json.loads(entry.info),
                'plan' : plan,
                'draft' : draft,
                'publish' : publish,
            }

            dataList.append( json.dumps(entryDict)  )

    #print('XXXXXX', dataList)
    random.shuffle(dataList)


    with open('static/topics.json', 'r') as f:
        srcJSON = json.load(f)


    sources = json.dumps(srcJSON['sources'])

    #print('DATA', type(dataList), dataList)
    return jsonify({'dataList' : dataList, 'sources' : sources, 'stage' : stage, 'info' : entryUser.info, 'partner': entryUser.partner, 'writer': entryUser.username })

## AJAX
@app.route('/getHTML/<string:unit>', methods=['POST'])
def getHTML(unit):
    print('GET HTML ACTIVE')

    try:
        name = request.form ['name']
        print('NAME', name, unit)
    except:
        name = current_user.username


    model = getInfo()['aModsDict'][unit]
    entry = model.query.filter_by(username=name).first()
    info = entry.info
    revise = entry.revise
    stage = entry.grade
    print(name, unit, stage, revise[0:5])

    #print('DATA', type(dataList), dataList)
    return jsonify({'revise' : revise, 'stage' : stage, 'info' : info})


""" ### PLAN/WORK/REVISE/PUBLISH ### """

@app.route("/work/<string:part>/<string:unit>", methods = ['GET', 'POST'])
@login_required
def part(part, unit):

    ### start project ###

    SCHEMA = getSchema()
    # S3_LOCATION = schemaList[SCHEMA]['S3_LOCATION']
    # S3_BUCKET_NAME = schemaList[SCHEMA]['S3_BUCKET_NAME']

    print('start project')

    classModel = getInfo()['aModsDict'][unit]

    entry = classModel.query.filter_by(partner=current_user.username).first()
    if not entry:
        entry = classModel.query.filter_by(username=current_user.username).first()

    if not entry:
        info = {
            'avatar' : 1,
            'theme' : 'white',
            'name' : current_user.username[0],
            'stage' : 0
        }
        # start assignment
        entry = classModel(username=current_user.username,
        info=json.dumps(info),
        plan=json.dumps({}),
        draft=json.dumps({}),
        revise=json.dumps({}),
        publish=json.dumps({})
        )
        db.session.add(entry)
        db.session.commit()

        entry = classModel.query.filter_by(username=current_user.username).first()

    fullDict = {
        'info' : entry.info,
        'plan' : entry.plan,
        'draft' : entry.draft,
        'revise' : entry.revise,
        'publish' : entry.publish,
    }

    #print(fullDict)

    #SOURCES = loadAWS('json_files/sources.json', 0)
    #sources = json.dumps(SOURCES['sources'])
    with open('static/topics.json', 'r') as f:
        srcJSON = json.load(f)

    sources = json.dumps(srcJSON['sources'])

    return render_template('work/' + part + '.html', unit=unit, fullDict=json.dumps(fullDict), sources=sources, partner=entry.partner)


""" ### INSTRUCTOR DASHBOARD ### """

@app.route("/write_dash", methods = ['GET', 'POST'])
@login_required
def write_dash():
    if current_user.id != 1:
        return redirect(url_for('home'))


    models_list = ['01', '02', '03', '04', '05']
    # models_list = ['05', '06', '07', '08', '09']
    recDict = {}

    userList = getWriteUsers()

    for model in getInfo()['aModsDict']:
        #print(recDict)
        if str(model) in models_list:
            recDict[model] = []
            for entry in getInfo()['aModsDict'][model].query.all():
                # if entry.partner == '':
                #     entry.partner = 'None'
                #     db.session.commit()
                info = json.loads(entry.info)
                user = entry.username
                vtm = 0

                try:
                    vtm = User.query.filter_by(username=user).first().vtm
                except:
                    print('VTM EXCEPT', user)

                entry = {
                    'user' : user,
                    'stage' : info['stage'],
                    'info' : info,
                    'partner' : entry.partner,
                    'plan' : json.loads(entry.plan),
                    'draft' : json.loads(entry.draft),
                    'revise' : json.loads(entry.revise),
                    'publish' : json.loads(entry.publish),
                    'vtm' : vtm
                }

                if len(recDict[model]) == 0:
                    recDict[model].append(entry)
                # else:
                #     recDict[model].append(entry)

                for e in recDict[model]:

                    print(e)
                    if int(info['stage']) < int(e['stage']):
                        recDict[model].insert(recDict[model].index(e), entry)
                        print(recDict[model].index(e))
                        break
                else:
                    recDict[model].append(entry)




    return  render_template('work/dashboard.html', recOBJ=str(json.dumps(recDict)))

@app.route("/published_work", methods = ['GET', 'POST'])
@login_required
def published():

    recDict = {}

    for model in getInfo()['aModsDict']:
        recDict[model] = {}
        #print(recDict)
        for entry in getInfo()['aModsDict'][model].query.all():
            recDict[str(model)][entry.username] = {
                'info' : json.loads(entry.info),
                'publish' : json.loads(entry.publish),
            }

    return  render_template('work/published_work.html', recOBJ=str(json.dumps(recDict)))



@app.route("/published_check", methods = ['POST'])
@login_required
def publish_API():

    recDict = {}

    for model in getInfo()['aModsDict']:
        recDict[model] = {}
        #print(recDict)
        for entry in getInfo()['aModsDict'][model].query.all():
            if entry.grade == 5 and int(model) > 4: # check models after topic 4
                reviseDict = json.loads(entry.revise)
                #print('xxxx', reviseDict)
                recDict[str(model)][entry.username] = {
                    'info' : json.loads(entry.info),
                    'publish' : json.loads(entry.publish),
                    'revise' : json.loads(entry.revise),
                    'plan' : json.loads(entry.plan),
                    'htmltext' : reviseDict['html'],
                }

    return jsonify({'revise' : revise, 'stage' : stage, 'info' : info})




@app.route("/published_check/<string:mode>", methods = ['GET', 'POST'])
@login_required
def pCheck(mode):
    taCheck = ['Chris', 'TA']

    recCount = 0
    recDict = {}

    for model in getInfo()['aModsDict']:
        recDict[model] = {}
        #print(recDict)
        for entry in getInfo()['aModsDict'][model].query.all():
            if current_user.username in taCheck:
                add = True
            elif current_user.username == entry.username:
                add = True
            else:
                add = False
            if entry.grade == 5 and add == True: # and int(model) > 4 -- check models after topic 4
                reviseDict = json.loads(entry.revise)
                #print('xxxx', reviseDict)
                recDict[str(model)][entry.username] = {
                    'info' : json.loads(entry.info),
                    'publish' : json.loads(entry.publish),
                    'revise' : json.loads(entry.revise),
                    'plan' : json.loads(entry.plan),
                    'htmltext' : reviseDict['html'],
                    'partner' : entry.partner
                }
                recCount += 1


    return  render_template('work/published_check.html', instructor=mode, recCount=recCount, recOBJ=str(json.dumps(recDict)))


@app.route("/editor/<string:student>/<string:unit>", methods = ['GET', 'POST'])
@login_required
def editor(student, unit):
    if current_user.id != 1:
        return redirect(url_for('home'))

    model = getInfo()['aModsDict'][unit]
    print(model)
    jStrings = model.query.filter_by(username=student).first()

    student_revise = jStrings.revise
    student_plan = jStrings.plan

    student_draft = json.loads(jStrings.draft)
    ## build the student text
    text = ''
    for part in student_draft:
        text += (student_draft[part] + ' ' )

    return  render_template('work/editor.html', text=text, student=student, unit=unit, student_revise=student_revise, student_plan=student_plan)




